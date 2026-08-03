import os
import sys
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as patches
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# --- COLOR PALETTE DEFINITIONS (Modern Dark Theme / Rust Inspired) ---
BG_DARK = RGBColor(26, 27, 32)        # #1A1B20 Charcoal Dark Slate
CARD_DARK = RGBColor(46, 49, 60)      # #2E313C Secondary Card Fill
BORDER_DARK = RGBColor(70, 75, 90)    # #464B5A Border Slate
RUST_ORANGE = RGBColor(229, 115, 36)  # #E57324 Primary Rust Orange
ACCENT_ORANGE = RGBColor(255, 138, 61)# #FF8A3D Bright Accent Orange
WHITE = RGBColor(243, 244, 246)       # #F3F4F6 Clean Off-White Text
MUTED_GREY = RGBColor(160, 166, 184)  # #A0A6B8 Muted Subtext Grey
GREEN_ACCENT = RGBColor(74, 222, 128) # #4ADE80 Success Accent

# --- MATPLOTLIB VISUAL ASSET GENERATION ---
def create_architecture_diagram(filename="arch_diagram.png"):
    fig, ax = plt.subplots(figsize=(11, 4.8), dpi=300)
    fig.patch.set_facecolor('#1A1B20')
    ax.set_facecolor('#1A1B20')
    ax.set_xlim(0, 10)
    ax.set_ylim(0, 5.5)
    ax.axis('off')

    # Main Pipeline Boxes
    boxes = [
        ("SRT File\n(.srt)", 0.3, 3.2, 1.4, 1.4, "#2E313C", "#E57324"),
        ("Parser\n\n• Zero-copy slices\n• encoding_rs BOM\n• WHATWG UTF-8", 2.2, 3.0, 1.8, 1.8, "#2E313C", "#FF8A3D"),
        ("Subtitle Data\nStructures\n\n• SubRipFile\n• SubRipItem\n• SubRipTime (i64)\n• Coordinates", 4.5, 2.9, 1.9, 2.0, "#24262E", "#E57324"),
        ("Editing APIs\n\n• Shift & Scale\n• Rate / FPS\n• Tag Stripping\n• CPS Calculator", 6.9, 3.0, 1.8, 1.8, "#2E313C", "#FF8A3D"),
        ("Writer\n\n• Deterministic EOL\n• BOM Handling\n• &str Serialize", 9.1, 3.2, 1.4, 1.4, "#2E313C", "#E57324")
    ]

    for title, x, y, w, h, bg, border in boxes:
        rect = patches.FancyBboxPatch((x - w/2, y - h/2), w, h, boxstyle="round,pad=0.1,rounding_size=0.15",
                                      facecolor=bg, edgecolor=border, linewidth=2.0)
        ax.add_patch(rect)
        ax.text(x, y, title, color="white", fontsize=10.5, fontweight="bold",
                ha="center", va="center", linespacing=1.35, family="sans-serif")

    # Arrows between boxes
    arrow_props = dict(facecolor="#FF8A3D", edgecolor="none", width=0.06, headwidth=6, headlength=7)
    arrow_pairs = [
        (1.0, 3.2, 1.3, 3.2),
        (3.1, 3.2, 3.55, 3.2),
        (5.45, 3.2, 6.0, 3.2),
        (7.8, 3.2, 8.4, 3.2)
    ]
    for x1, y1, x2, y2 in arrow_pairs:
        ax.annotate('', xy=(x2, y2), xytext=(x1, y1), arrowprops=arrow_props)

    # Export Interfaces (Bottom Layer)
    iface_boxes = [
        ("libsrt (PyO3 Python Extension Bridge — #![feature(python)])\n• 100% Unmodified Pytest Compatibility  • Python Magic Methods  • Explicit ErrorHandling Enum", 
         3.3, 0.85, 4.4, 1.1, "#20222B", "#4ADE80"),
        ("srt (Native Standalone Static CLI Binary)\n• Built with clap  • <1ms Cold Start  • Zero Python Runtime Dependency", 
         7.8, 0.85, 3.6, 1.1, "#20222B", "#60A5FA")
    ]
    for text, x, y, w, h, bg, border in iface_boxes:
        rect = patches.FancyBboxPatch((x - w/2, y - h/2), w, h, boxstyle="round,pad=0.1,rounding_size=0.15",
                                      facecolor=bg, edgecolor=border, linewidth=1.8, linestyle="--")
        ax.add_patch(rect)
        ax.text(x, y, text, color="#F3F4F6", fontsize=9.5, fontweight="semibold",
                ha="center", va="center", linespacing=1.3, family="sans-serif")

    # Connecting vertical arrows to interfaces
    ax.annotate('', xy=(3.3, 1.4), xytext=(4.0, 1.9),
                arrowprops=dict(facecolor="#4ADE80", edgecolor="none", width=0.04, headwidth=5, headlength=6))
    ax.annotate('', xy=(7.2, 1.4), xytext=(6.9, 2.1),
                arrowprops=dict(facecolor="#60A5FA", edgecolor="none", width=0.04, headwidth=5, headlength=6))

    plt.tight_layout()
    plt.savefig(filename, dpi=300, bbox_inches='tight', facecolor='#1A1B20')
    plt.close()
    print(f"Generated {filename}")

def create_benchmark_chart(filename="benchmark_chart.png"):
    fig, ax = plt.subplots(figsize=(8.5, 4.2), dpi=300)
    fig.patch.set_facecolor('#1A1B20')
    ax.set_facecolor('#1A1B20')

    ops = [
        "Serialize (text)\n1000 items",
        "Shift Timestamps\n1000 items",
        "Parse (Movie - 1332)\nutf-8.srt",
        "Parse (1000 subs)\nsynthetic"
    ]
    py_times = [0.15, 2.11, 9.05, 6.58]
    rs_times = [0.04, 0.44, 0.50, 0.30]
    speedups = ["3.4x Faster", "4.8x Faster", "18.0x Faster", "22.0x Faster"]

    y = range(len(ops))
    height = 0.34

    bars_py = ax.barh([i + height/2 for i in y], py_times, height=height, color='#464B5A', label='Python (byroot/pysrt)', edgecolor='#70758A')
    bars_rs = ax.barh([i - height/2 for i in y], rs_times, height=height, color='#FF8A3D', label='Rust (pysrt-rs)', edgecolor='#E57324', linewidth=1.5)

    ax.set_xlabel('Execution Time in Milliseconds (Lower is Better)', color='#F3F4F6', fontsize=11, fontweight='bold', labelpad=10)
    ax.set_yticks(y)
    ax.set_yticklabels(ops, color='#F3F4F6', fontsize=10.5, fontweight='semibold')
    ax.tick_params(axis='x', colors='#A0A6B8', labelsize=10)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['left'].set_color('#464B5A')
    ax.spines['bottom'].set_color('#464B5A')
    ax.xaxis.grid(True, linestyle='--', color='#2E313C', alpha=0.7)
    ax.set_axisbelow(True)

    for idx, (py_val, rs_val, sp) in enumerate(zip(py_times, rs_times, speedups)):
        ax.text(py_val + 0.25, idx + height/2, f"{py_val:.2f} ms", va='center', color='#A0A6B8', fontsize=9.5)
        ax.text(rs_val + 0.25, idx - height/2, f"{rs_val:.2f} ms  [{sp}]", va='center', color='#FF8A3D', fontsize=10, fontweight='bold')

    ax.set_xlim(0, 12.5)
    legend = ax.legend(loc='upper right', frameon=True, facecolor='#2E313C', edgecolor='#464B5A', fontsize=10)
    for text in legend.get_texts():
        text.set_color('#F3F4F6')

    plt.title('Performance Speedup: Rust Port vs Python Upstream (Median of 5 Rounds, 1000 Iterations)', 
              color='#F3F4F6', fontsize=12, fontweight='bold', pad=15)
    plt.tight_layout()
    plt.savefig(filename, dpi=300, bbox_inches='tight', facecolor='#1A1B20')
    plt.close()
    print(f"Generated {filename}")

def create_workflow_card(filename="workflow_card.png"):
    fig, ax = plt.subplots(figsize=(10.5, 4.3), dpi=300)
    fig.patch.set_facecolor('#1A1B20')
    ax.set_facecolor('#1A1B20')
    ax.set_xlim(0, 10)
    ax.set_ylim(0, 5)
    ax.axis('off')

    # Left Card: Idiomatic Rust Usage
    rect_left = patches.FancyBboxPatch((0.2, 0.3), 4.6, 4.4, boxstyle="round,pad=0.1,rounding_size=0.15",
                                       facecolor="#24262E", edgecolor="#E57324", linewidth=1.5)
    ax.add_patch(rect_left)
    ax.text(0.5, 4.35, "Idiomatic Rust Core Usage (src/lib.rs)", color="#FF8A3D", fontsize=11, fontweight="bold")
    
    code_text = (
        "use libsrt::{SubRipFile, SubRipTime};\n\n"
        "// 1. Zero-copy parse & BOM detection\n"
        "let mut file = SubRipFile::open(\"movie.srt\")?;\n\n"
        "// 2. Safe timestamp shift across all items\n"
        "for item in &mut file.items {\n"
        "    item.shift(SubRipTime::from_ms(500));\n"
        "}\n\n"
        "// 3. Deterministic EOL serialization\n"
        "file.save(\"movie_shifted.srt\")?;"
    )
    ax.text(0.5, 2.3, code_text, color="#F3F4F6", fontsize=9.8, family="monospace", linespacing=1.4)

    # Right Card: Repository Architecture & Verification
    rect_right = patches.FancyBboxPatch((5.2, 0.3), 4.6, 4.4, boxstyle="round,pad=0.1,rounding_size=0.15",
                                        facecolor="#24262E", edgecolor="#4ADE80", linewidth=1.5)
    ax.add_patch(rect_right)
    ax.text(5.5, 4.35, "Repository Structure & Verification", color="#4ADE80", fontsize=11, fontweight="bold")

    tree_text = (
        "pysrt-rs/\n"
        " ├── src/\n"
        " │    ├── lib.rs     (#![forbid(unsafe_code)])\n"
        " │    ├── time.rs    (Normalized i64 Ordinal)\n"
        " │    ├── item.rs    (Zero-Copy Parse & CPS)\n"
        " │    ├── file.rs    (encoding_rs & BOM Sniff)\n"
        " │    └── bin/srt.rs (Standalone Native CLI)\n"
        " ├── fuzz/diff_fuzz.py  (7,000 cases, 0 diffs)\n"
        " ├── tests/port/        (75 native Rust tests)\n"
        " ├── tests/fixed/       (75 Python tests pass)\n"
        " └── DECISIONS.md       (10 Arch Decisions)"
    )
    ax.text(5.5, 2.1, tree_text, color="#F3F4F6", fontsize=9.5, family="monospace", linespacing=1.35)

    plt.tight_layout()
    plt.savefig(filename, dpi=300, bbox_inches='tight', facecolor='#1A1B20')
    plt.close()
    print(f"Generated {filename}")

# --- PPTX HELPER FUNCTIONS ---
def add_dark_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_DARK

def add_header(slide, title_text, category_text="PORT MORTEM 2026 • TRACK D: PYTHON → RUST"):
    # Accent top colored line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.4), Inches(11.73), Inches(0.05)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = RUST_ORANGE
    line.line.color.rgb = RUST_ORANGE

    # Category tracker / Eyebrow
    cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.48), Inches(11.0), Inches(0.3))
    tf_cat = cat_box.text_frame
    tf_cat.word_wrap = True
    p_cat = tf_cat.paragraphs[0]
    p_cat.text = category_text.upper()
    p_cat.font.size = Pt(9.5)
    p_cat.font.bold = True
    p_cat.font.color.rgb = ACCENT_ORANGE
    p_cat.font.name = "Arial"

    # Main Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.72), Inches(11.5), Inches(0.75))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    p_title = tf_title.paragraphs[0]
    p_title.text = title_text
    p_title.font.size = Pt(28)
    p_title.font.bold = True
    p_title.font.color.rgb = WHITE
    p_title.font.name = "Arial"

def add_card(slide, left, top, width, height, bg_color=CARD_DARK, border_color=BORDER_DARK):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    card.fill.solid()
    card.fill.fore_color.rgb = bg_color
    card.line.color.rgb = border_color
    card.line.width = Pt(1.5)
    return card

def add_bullet_list(tf, items, font_size=15, line_spacing=1.2, color=WHITE):
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Arial"
        p.space_after = Pt(8)
        p.line_spacing = line_spacing

def add_footer(slide, current_slide, total_slides=9):
    footer_box = slide.shapes.add_textbox(Inches(0.8), Inches(7.0), Inches(11.73), Inches(0.35))
    tf = footer_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"pysrt-rs: Memory-Safe Rust Port of byroot/pysrt    |    Slide {current_slide} of {total_slides}"
    p.font.size = Pt(10)
    p.font.color.rgb = MUTED_GREY
    p.font.name = "Arial"

# --- MAIN PRESENTATION BUILDER ---
def build_presentation(filename="pysrt-rs_Port_Mortem_2026.pptx"):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # ==========================================
    # SLIDE 1: TITLE SLIDE
    # ==========================================
    s1 = prs.slides.add_slide(blank_layout)
    add_dark_background(s1)

    # Decorative left accent block
    acc = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.8), Inches(0.18), Inches(4.0))
    acc.fill.solid()
    acc.fill.fore_color.rgb = RUST_ORANGE
    acc.line.color.rgb = RUST_ORANGE

    # Title card container
    tbox = s1.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(11.0), Inches(4.2))
    tf1 = tbox.text_frame
    tf1.word_wrap = True

    p0 = tf1.paragraphs[0]
    p0.text = "PORT MORTEM 2026  •  HACKATHON SUBMISSION  •  TRACK D"
    p0.font.size = Pt(13)
    p0.font.bold = True
    p0.font.color.rgb = ACCENT_ORANGE
    p0.font.name = "Arial"
    p0.space_after = Pt(15)

    p1 = tf1.add_paragraph()
    p1.text = "pysrt-rs"
    p1.font.size = Pt(54)
    p1.font.bold = True
    p1.font.color.rgb = WHITE
    p1.font.name = "Arial"
    p1.space_after = Pt(10)

    p2 = tf1.add_paragraph()
    p2.text = "Porting pysrt from Python to High-Performance, Memory-Safe Rust"
    p2.font.size = Pt(22)
    p2.font.color.rgb = MUTED_GREY
    p2.font.name = "Arial"
    p2.space_after = Pt(40)

    p3 = tf1.add_paragraph()
    p3.text = "Author: K. Anirudh Reddy"
    p3.font.size = Pt(18)
    p3.font.bold = True
    p3.font.color.rgb = WHITE
    p3.font.name = "Arial"
    p3.space_after = Pt(6)

    p4 = tf1.add_paragraph()
    p4.text = "Repository: https://github.com/Anirudh-12/pysrt-rs   |   Upstream: https://github.com/byroot/pysrt"
    p4.font.size = Pt(14)
    p4.font.color.rgb = ACCENT_ORANGE
    p4.font.name = "Arial"

    # ==========================================
    # SLIDE 2: PROJECT OVERVIEW
    # ==========================================
    s2 = prs.slides.add_slide(blank_layout)
    add_dark_background(s2)
    add_header(s2, "Project Overview: Why Port pysrt to Rust?")
    add_footer(s2, 2)

    # 3 Structured Columns
    cols = [
        ("WHAT IS PYSRT?", 0.8, 1.7, 3.65, 4.9, [
            "• Industry-Standard Python Library: Created by @byroot (byroot/pysrt) for SubRip (.srt) parsing, editing, and serialization.",
            "• Core Subtitle Workflows: Manages timestamps, subtitle blocks, text formatting tags, positioning coordinates, and time-shifting.",
            "• Standalone Utility: Widely deployed in video encoding pipelines, NLP datasets, and accessibility captioning automation."
        ]),
        ("WHAT PROBLEMS IT SOLVES", 4.84, 1.7, 3.65, 4.9, [
            "• Heterogeneous Encodings: Automatically decodes UTF-8, UTF-8-sig (BOM), UTF-16 LE/BE, and legacy encodings.",
            "• Temporal Synchronization: Precise millisecond time shifting, frame-rate (FPS) conversions, and duration calculations.",
            "• EOL & Formatting Consistency: Normalizes CRLF / LF line endings and calculates reading speed (Characters Per Second - CPS)."
        ]),
        ("WHY RUST IS THE TARGET", 8.88, 1.7, 3.65, 4.9, [
            "• Guaranteed Memory Safety: #![forbid(unsafe_code)] in core crate; zero GC pauses or heap fragmentation.",
            "• Zero-Copy Parsing: High-speed str / byte slice scanners replace expensive Python regular expression matching.",
            "• Dual Architecture: Exports both an instant-startup CLI binary (srt) and a Python extension bridge (libsrt via PyO3)."
        ])
    ]

    for title, left, top, w, h, bullets in cols:
        add_card(s2, left, top, w, h)
        box = s2.shapes.add_textbox(Inches(left + 0.2), Inches(top + 0.2), Inches(w - 0.4), Inches(h - 0.4))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = ACCENT_ORANGE
        p.space_after = Pt(14)
        for b in bullets:
            p = tf.add_paragraph()
            p.text = b
            p.font.size = Pt(13)
            p.font.color.rgb = WHITE
            p.space_after = Pt(12)

    # ==========================================
    # SLIDE 3: PORTING GOALS
    # ==========================================
    s3 = prs.slides.add_slide(blank_layout)
    add_dark_background(s3)
    add_header(s3, "Engineering Goals for the Port")
    add_footer(s3, 3)

    grid = [
        ("1. 100% Behavioral Equivalence", 0.8, 1.7, 5.65, 2.3, [
            "• Bug-for-bug parity with original Python semantics across 75 integration tests.",
            "• Verified against 7,000 differential fuzzing cases (fuzz/diff_fuzz.py) with 0 divergences.",
            "• Identical string representation, timestamp arithmetic, and tag handling."
        ]),
        ("2. Guaranteed Memory Safety & #![forbid(unsafe_code)]", 6.88, 1.7, 5.65, 2.3, [
            "• Core library crate (src/lib.rs, time.rs, item.rs, file.rs) contains zero unsafe blocks.",
            "• Eliminates undefined behavior, null pointer dereferences, and buffer overflows.",
            "• Earned the +5 Zero Unsafe bonus in Port Mortem hackathon scoring."
        ]),
        ("3. Strong Typing & Invariant Protection", 0.8, 4.25, 5.65, 2.3, [
            "• Normalized scalar integer (i64 ms ordinal) replaces mutable 4-field datetime objects.",
            "• Strongly typed Coordinates struct replaces unstructured Python X1:Y1 dictionaries.",
            "• Monadic Result<T, SrtError> makes failure modes explicit in the type system."
        ]),
        ("4. Architectural Maintainability & Performance", 6.88, 4.25, 5.65, 2.3, [
            "• Clear separation of concerns: core parsing, CLI binary, and optional PyO3 FFI bridge.",
            "• 10 architectural decision records documented in DECISIONS.md.",
            "• Zero-copy byte/str slice scanning delivers up to 22.0× parse speedup."
        ])
    ]

    for title, left, top, w, h, bullets in grid:
        add_card(s3, left, top, w, h)
        box = s3.shapes.add_textbox(Inches(left + 0.25), Inches(top + 0.2), Inches(w - 0.5), Inches(h - 0.4))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = ACCENT_ORANGE
        p.space_after = Pt(10)
        for b in bullets:
            p = tf.add_paragraph()
            p.text = b
            p.font.size = Pt(13)
            p.font.color.rgb = WHITE
            p.space_after = Pt(6)

    # ==========================================
    # SLIDE 4: ARCHITECTURE
    # ==========================================
    s4 = prs.slides.add_slide(blank_layout)
    add_dark_background(s4)
    add_header(s4, "System Architecture: SubRip Pipeline & Dual Interfaces")
    add_footer(s4, 4)

    # Insert Architecture Diagram Image
    s4.shapes.add_picture("arch_diagram.png", Inches(0.8), Inches(1.6), Inches(11.73), Inches(5.1))

    # ==========================================
    # SLIDE 5: ENGINEERING CHALLENGES
    # ==========================================
    s5 = prs.slides.add_slide(blank_layout)
    add_dark_background(s5)
    add_header(s5, "Technical Challenges & Architectural Decisions (DECISIONS.md)")
    add_footer(s5, 5)

    challenges = [
        ("Dynamic Typing → Static Typing (DECISIONS.md §2, §9)", 0.8, 1.7, 5.65, 2.3, [
            "• Challenge: Upstream used unstructured dictionaries for subtitle positioning coordinates (X1: ... Y1: ...) and stringly-typed error modes ('ERROR_PASS', 'ERROR_LOG').",
            "• Solution: Built a strongly typed Coordinates struct and an explicit ErrorHandling enum. Python error modes are cleanly mapped at the PyO3 FFI boundary."
        ]),
        ("Time Representation & Invariants (DECISIONS.md §1)", 6.88, 1.7, 5.65, 2.3, [
            "• Challenge: Upstream SubRipTime stored mutable hours, minutes, seconds, and milliseconds, allowing invalid states (e.g., negative ms or seconds > 59).",
            "• Solution: Stored time as a single normalized scalar integer (i64 total milliseconds ordinal). Property accessors are computed on demand; arithmetic is O(1)."
        ]),
        ("Zero-Copy Parsing & Encoding Edge Cases (DECISIONS.md §3, §4, §8)", 0.8, 4.25, 5.65, 2.3, [
            "• Challenge: Python relied on regular expressions (re) and slow heuristic character detection (chardet) with unpredictable EOL behavior.",
            "• Solution: Built custom zero-copy slice scanners without regex engines. Integrated Mozilla's encoding_rs for SIMD-accelerated, WHATWG-compliant BOM stripping."
        ]),
        ("Ownership, Borrowing & Zero Unsafe (DECISIONS.md §5, §7)", 6.88, 4.25, 5.65, 2.3, [
            "• Challenge: Porting a GC-managed Python object graph to safe Rust without reference cycles while maintaining unmodified Python test compatibility.",
            "• Solution: SubRipItem owns text (String) while SubRipTime is Copy. Isolated all FFI glue inside PyO3 macros (src/python/mod.rs); 0 unsafe blocks in core."
        ])
    ]

    for title, left, top, w, h, bullets in challenges:
        add_card(s5, left, top, w, h)
        box = s5.shapes.add_textbox(Inches(left + 0.25), Inches(top + 0.2), Inches(w - 0.5), Inches(h - 0.4))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = ACCENT_ORANGE
        p.space_after = Pt(10)
        for b in bullets:
            p = tf.add_paragraph()
            p.text = b
            p.font.size = Pt(13)
            p.font.color.rgb = WHITE
            p.space_after = Pt(8)

    # ==========================================
    # SLIDE 6: TESTING & CORRECTNESS (IMPORTANT SLIDE)
    # ==========================================
    s6 = prs.slides.add_slide(blank_layout)
    add_dark_background(s6)
    add_header(s6, "Validation, Test Suite Fidelity & Correctness Note")
    add_footer(s6, 6)

    # Top Quote Card (Word-for-Word Required Engineering Note)
    quote_card = add_card(s6, 0.8, 1.7, 11.73, 1.6, bg_color=RGBColor(38, 41, 52), border_color=ACCENT_ORANGE)
    qbox = s6.shapes.add_textbox(Inches(1.0), Inches(1.85), Inches(11.3), Inches(1.3))
    qtf = qbox.text_frame
    qtf.word_wrap = True
    qp0 = qtf.paragraphs[0]
    qp0.text = "CRITICAL ENGINEERING VALIDATION NOTE:"
    qp0.font.size = Pt(12)
    qp0.font.bold = True
    qp0.font.color.rgb = ACCENT_ORANGE
    qp0.space_after = Pt(6)
    qp1 = qtf.add_paragraph()
    qp1.text = (
        "\"During validation, an issue was discovered in the upstream project's test suite. The original upstream "
        "repository also fails this particular test because of a bug in the test itself. Rather than reproducing an "
        "incorrect test expectation, the test was corrected to reflect the intended behavior while preserving the "
        "library's functionality.\""
    )
    qp1.font.size = Pt(13.5)
    qp1.font.italic = True
    qp1.font.color.rgb = WHITE
    qp1.line_spacing = 1.25

    # Bottom Two Explanation Columns
    add_card(s6, 0.8, 3.55, 5.65, 3.1)
    tbox_l = s6.shapes.add_textbox(Inches(1.0), Inches(3.75), Inches(5.25), Inches(2.7))
    tfl = tbox_l.text_frame
    tfl.word_wrap = True
    pl0 = tfl.paragraphs[0]
    pl0.text = "Why test_save Fails in Unmodified Upstream"
    pl0.font.size = Pt(15)
    pl0.font.bold = True
    pl0.font.color.rgb = ACCENT_ORANGE
    pl0.space_after = Pt(10)
    add_bullet_list(tfl, [
        "• Specific Test File: tests/test_srtfile.py -> TestSerialization::test_save.",
        "• Root Cause of Assertion Failure: test_save saves a file with eol='\\n' and compares byte-for-byte against static reference fixture tests/static/utf-8.srt.",
        "• CRLF Mismatch: The static fixture utf-8.srt in upstream was committed with Windows CRLF ('\\r\\n'), making b'0\\n...' == b'0\\r\\n...' impossible.",
        "• Important Clarification: The issue is purely with the test fixture expectation, NOT the library implementation. The unmodified upstream Python library fails this exact test under Python 3."
    ], font_size=12.5, line_spacing=1.2)

    add_card(s6, 6.88, 3.55, 5.65, 3.1)
    tbox_r = s6.shapes.add_textbox(Inches(7.08), Inches(3.75), Inches(5.25), Inches(2.7))
    tfr = tbox_r.text_frame
    tfr.word_wrap = True
    pr0 = tfr.paragraphs[0]
    pr0.text = "Three-Layer Verification Evidence"
    pr0.font.size = Pt(15)
    pr0.font.bold = True
    pr0.font.color.rgb = ACCENT_ORANGE
    pr0.space_after = Pt(10)
    add_bullet_list(tfr, [
        "• Unmodified Upstream Suite (tests/original/): Exactly 74 / 75 tests pass (test_srttime 21/21, test_srtitem 18/18, test_srtfile 35/36 — 1 fixture bug).",
        "• Corrected Suite (tests/fixed/): 75 / 75 tests pass (100% parity) where test_save is corrected to use eol='\\r\\n'. Matches upstream reference Python 100%.",
        "• Native Rust Suite (tests/port/): 100% native Rust integration suite ports all 75 tests 1-to-1 (75 / 75 pass) + 9 unit tests across lib.rs & CLI (84 / 84 total pass).",
        "• Zero File Modifications in Original: File SHA-256 hashes are strictly verified in .port-mortem.toml."
    ], font_size=12.5, line_spacing=1.2)

    # ==========================================
    # SLIDE 7: RESULTS & PERFORMANCE BENCHMARKS
    # ==========================================
    s7 = prs.slides.add_slide(blank_layout)
    add_dark_background(s7)
    add_header(s7, "Empirical Results: Performance, Parity & Repository Workflow")
    add_footer(s7, 7)

    # Left Side: Benchmark Chart Image
    s7.shapes.add_picture("benchmark_chart.png", Inches(0.8), Inches(1.6), Inches(6.5), Inches(3.3))

    # Left Bottom Table / Metric Badge
    metric_card = add_card(s7, 0.8, 5.05, 6.5, 1.65, bg_color=RGBColor(36, 38, 46))
    mbox = s7.shapes.add_textbox(Inches(1.0), Inches(5.15), Inches(6.1), Inches(1.45))
    mtf = mbox.text_frame
    mtf.word_wrap = True
    mp0 = mtf.paragraphs[0]
    mp0.text = "BENCHMARK METHODOLOGY & HIGHLIGHTS (bench/run_bench.py)"
    mp0.font.size = Pt(12)
    mp0.font.bold = True
    mp0.font.color.rgb = ACCENT_ORANGE
    mp0.space_after = Pt(6)
    add_bullet_list(mtf, [
        "• Up to 22.0× Throughput Speedup: Parses 1,000 subtitles in 0.30 ms vs 6.58 ms for Python.",
        "• 4.8× Faster Timestamp Shifting & 3.4× Faster Serialization across 1,332-item movie subtitles.",
        "• 7,000 / 7,000 Differential Fuzzing Cases Pass: Zero divergence (fuzz/diff_fuzz.py) vs Python.",
        "• 100% Core Build Success: #![forbid(unsafe_code)] enforced across crate root."
    ], font_size=11.5, line_spacing=1.15)

    # Right Side: Workflow Code & Repo Structure Image
    s7.shapes.add_picture("workflow_card.png", Inches(7.5), Inches(1.6), Inches(5.03), Inches(5.1))

    # ==========================================
    # SLIDE 8: LESSONS LEARNED
    # ==========================================
    s8 = prs.slides.add_slide(blank_layout)
    add_dark_background(s8)
    add_header(s8, "Lessons Learned & Engineering Insights")
    add_footer(s8, 8)

    lessons = [
        ("1. Data Representation is Everything", 0.8, 1.7, 5.65, 2.3, [
            "• Replacing a 4-field mutable datetime object with an 8-byte scalar integer (i64 ms ordinal) eliminated entire classes of invalid time states.",
            "• Scalar integer arithmetic made comparison and timestamp arithmetic instantaneous and O(1)."
        ]),
        ("2. Designing Clean FFI Boundaries", 6.88, 1.7, 5.65, 2.3, [
            "• Maintaining 100% bug-for-bug compatibility for Python error recovery modes ('ERROR_PASS') required careful separation of concerns.",
            "• Explicit ErrorHandling enum and callbacks at the PyO3 layer prevented Python semantics from polluting idiomatic Rust core logic."
        ]),
        ("3. WHATWG Encoding Standards vs Heuristics", 0.8, 4.25, 5.65, 2.3, [
            "• Standardized WHATWG encoding detection (encoding_rs) is vastly superior to heuristic-based character sniffers like chardet.",
            "• Handles multi-byte BOM stripping (UTF-8-sig, UTF-16, UTF-32) deterministically without runtime interpreter locks."
        ]),
        ("4. The Power of Differential Fuzzing", 6.88, 4.25, 5.65, 2.3, [
            "• Running 7,000 randomized test cases simultaneously through Rust and Python proved indispensable for verifying behavioral equivalence.",
            "• Fuzzing even caught a latent timestamp-overflow bug in the upstream test generator (00:00:99,500) where seconds exceeded 59!"
        ])
    ]

    for title, left, top, w, h, bullets in lessons:
        add_card(s8, left, top, w, h)
        box = s8.shapes.add_textbox(Inches(left + 0.25), Inches(top + 0.2), Inches(w - 0.5), Inches(h - 0.4))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = ACCENT_ORANGE
        p.space_after = Pt(10)
        for b in bullets:
            p = tf.add_paragraph()
            p.text = b
            p.font.size = Pt(13)
            p.font.color.rgb = WHITE
            p.space_after = Pt(8)

    # ==========================================
    # SLIDE 9: RESOURCES & LINKS
    # ==========================================
    s9 = prs.slides.add_slide(blank_layout)
    add_dark_background(s9)
    add_header(s9, "Project Resources, Documentation & Submission Links")
    add_footer(s9, 9)

    add_card(s9, 0.8, 1.7, 11.73, 4.0)
    rbox = s9.shapes.add_textbox(Inches(1.2), Inches(1.95), Inches(10.9), Inches(3.5))
    rtf = rbox.text_frame
    rtf.word_wrap = True

    rp0 = rtf.paragraphs[0]
    rp0.text = "PORT MORTEM 2026 — TRACK D SUBMISSION REFERENCES"
    rp0.font.size = Pt(16)
    rp0.font.bold = True
    rp0.font.color.rgb = ACCENT_ORANGE
    rp0.space_after = Pt(16)

    links = [
        ("pysrt-rs Repository (Rust Port)", "https://github.com/Anirudh-12/pysrt-rs", "Complete source code, CLI binary (src/bin/srt.rs), and test suites."),
        ("Original Upstream Repository", "https://github.com/byroot/pysrt", "Original pure-Python SubRip library by @byroot."),
        ("Standalone Rust Library Docs", "RUST_LIBRARY.md (in repository)", "Detailed guide for using libsrt as a standalone Rust crate (cdylib & rlib)."),
        ("Architectural Decision Records", "DECISIONS.md (in repository)", "10 comprehensive architectural records detailing Python -> Rust trade-offs."),
        ("Port Mortem Hackathon Manifest", ".port-mortem.toml (in repository)", "Submission metadata, SHA-256 file verification hashes, and bonus claims.")
    ]

    for name, url, desc in links:
        p = rtf.add_paragraph()
        p.text = f"• {name}:  "
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.space_after = Pt(2)
        
        # Add URL sub-line
        p_url = rtf.add_paragraph()
        p_url.text = f"    {url} — {desc}"
        p_url.font.size = Pt(12.5)
        p_url.font.color.rgb = ACCENT_ORANGE
        p_url.space_after = Pt(12)

    # Bottom Summary Badge
    badge_card = add_card(s9, 0.8, 5.9, 11.73, 0.85, bg_color=RGBColor(38, 41, 52), border_color=ACCENT_ORANGE)
    bbox = s9.shapes.add_textbox(Inches(1.0), Inches(6.08), Inches(11.3), Inches(0.5))
    btf = bbox.text_frame
    bp0 = btf.paragraphs[0]
    bp0.text = "100% NATIVE RUST PARITY   •   22× FASTER PARSING   •   0 UNSAFE CORE   •   7,000 DIFF FUZZ SURVIVOR"
    bp0.font.size = Pt(13)
    bp0.font.bold = True
    bp0.font.color.rgb = WHITE
    bp0.alignment = PP_ALIGN.CENTER

    # Save presentation
    prs.save(filename)
    print(f"Presentation saved successfully to: {filename}")

if __name__ == "__main__":
    create_architecture_diagram()
    create_benchmark_chart()
    create_workflow_card()
    build_presentation()
